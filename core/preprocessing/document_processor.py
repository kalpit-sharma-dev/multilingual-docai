"""
Multi-format Document Processor for PS-05

Converts various document formats (PDF, DOC, PPT, etc.) to images
for layout detection processing. Supports offline deployment.
"""

import logging
import os
from pathlib import Path
from typing import List, Tuple, Optional, Dict
import cv2
import numpy as np

logger = logging.getLogger(__name__)

class DocumentProcessor:
    """Multi-format document processor for PS-05."""
    
    def __init__(self, config: Dict = None):
        """Initialize document processor.
        
        Args:
            config: Configuration dictionary
        """
        self.config = config or {}
        self.supported_formats = {
            '.pdf': self._process_pdf,
            '.doc': self._process_doc,
            '.docx': self._process_docx,
            '.ppt': self._process_ppt,
            '.pptx': self._process_pptx,
            '.png': self._process_image,
            '.jpg': self._process_image,
            '.jpeg': self._process_image,
            '.bmp': self._process_image,
            '.tiff': self._process_image,
            '.tif': self._process_image
        }
        
        # Check for required dependencies
        self._check_dependencies()
    
    def _check_dependencies(self):
        """Check if required dependencies are available."""
        self.pdf_available = False
        self.doc_available = False
        self.ppt_available = False
        
        # Check PDF support
        try:
            import fitz  # PyMuPDF
            self.pdf_available = True
            logger.info("PDF processing available (PyMuPDF)")
        except ImportError:
            try:
                import pdf2image
                self.pdf_available = True
                logger.info("PDF processing available (pdf2image)")
            except ImportError:
                logger.warning("PDF processing not available. Install PyMuPDF or pdf2image")
        
        # Check DOC support
        try:
            import docx2txt
            self.doc_available = True
            logger.info("DOC processing available (docx2txt)")
        except ImportError:
            logger.warning("DOC processing not available. Install docx2txt")
        
        # Check PPT support
        try:
            import python_pptx
            self.ppt_available = True
            logger.info("PPT processing available (python-pptx)")
        except ImportError:
            logger.warning("PPT processing not available. Install python-pptx")
    
    def can_process(self, file_path: str) -> bool:
        """Check if file can be processed.
        
        Args:
            file_path: Path to file
            
        Returns:
            True if file can be processed
        """
        suffix = Path(file_path).suffix.lower()
        return suffix in self.supported_formats
    
    def get_supported_formats(self) -> List[str]:
        """Get list of supported file formats.
        
        Returns:
            List of supported file extensions
        """
        return list(self.supported_formats.keys())
    
    def process_document(self, file_path: str, output_dir: str = None) -> List[str]:
        """Process document and convert to images.
        
        Args:
            file_path: Path to input document
            output_dir: Output directory for images (optional)
            
        Returns:
            List of generated image paths
        """
        file_path = Path(file_path)
        if not file_path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")
        
        suffix = file_path.suffix.lower()
        if suffix not in self.supported_formats:
            raise ValueError(f"Unsupported file format: {suffix}")
        
        # Create output directory
        if output_dir is None:
            output_dir = file_path.parent / f"{file_path.stem}_images"
        else:
            output_dir = Path(output_dir) / f"{file_path.stem}_images"
        
        output_dir.mkdir(parents=True, exist_ok=True)
        
        # Process document
        processor = self.supported_formats[suffix]
        try:
            image_paths = processor(str(file_path), str(output_dir))
            logger.info(f"Processed {file_path} -> {len(image_paths)} images")
            return image_paths
        except Exception as e:
            logger.error(f"Error processing {file_path}: {e}")
            raise
    
    def _process_pdf(self, file_path: str, output_dir: str) -> List[str]:
        """Process PDF document.
        
        Args:
            file_path: Path to PDF file
            output_dir: Output directory for images
            
        Returns:
            List of generated image paths
        """
        if not self.pdf_available:
            raise RuntimeError("PDF processing not available")
        
        try:
            # Try PyMuPDF first
            import fitz
            return self._process_pdf_pymupdf(file_path, output_dir)
        except ImportError:
            # Fallback to pdf2image
            try:
                import pdf2image
                return self._process_pdf_pdf2image(file_path, output_dir)
            except ImportError:
                raise RuntimeError("No PDF processing library available")
    
    def _process_pdf_pymupdf(self, file_path: str, output_dir: str) -> List[str]:
        """Process PDF using PyMuPDF."""
        import fitz
        
        doc = fitz.open(file_path)
        image_paths = []
        
        for page_num in range(len(doc)):
            page = doc.load_page(page_num)
            
            # Set zoom factor for better quality
            mat = fitz.Matrix(2.0, 2.0)  # 2x zoom
            pix = page.get_pixmap(matrix=mat)
            
            # Convert to PIL Image
            img_data = pix.tobytes("png")
            img = cv2.imdecode(np.frombuffer(img_data, np.uint8), cv2.IMREAD_COLOR)
            
            # Save image
            output_path = Path(output_dir) / f"page_{page_num:03d}.png"
            cv2.imwrite(str(output_path), img)
            image_paths.append(str(output_path))
        
        doc.close()
        return image_paths
    
    def _process_pdf_pdf2image(self, file_path: str, output_dir: str) -> List[str]:
        """Process PDF using pdf2image."""
        import pdf2image
        
        # Convert PDF to images
        images = pdf2image.convert_from_path(file_path, dpi=300)
        image_paths = []
        
        for i, image in enumerate(images):
            output_path = Path(output_dir) / f"page_{i:03d}.png"
            image.save(str(output_path), "PNG")
            image_paths.append(str(output_path))
        
        return image_paths
    
    def _process_doc(self, file_path: str, output_dir: str) -> List[str]:
        """Process DOC document.
        
        Args:
            file_path: Path to DOC file
            output_dir: Output directory for images
            
        Returns:
            List of generated image paths
        """
        if not self.doc_available:
            raise RuntimeError("DOC processing not available")
        
        # For now, convert to DOCX first, then process
        # This is a simplified approach - in production you might want to use
        # more sophisticated DOC processing libraries
        logger.warning("DOC processing is limited. Consider converting to DOCX first.")
        
        # Create a placeholder image with text content
        output_path = Path(output_dir) / "page_000.png"
        
        # Create a simple image with document info
        img = np.ones((800, 600, 3), dtype=np.uint8) * 255
        cv2.putText(img, f"DOC Document: {Path(file_path).name}", (50, 100), 
                   cv2.FONT_HERSHEY_SIMPLEX, 1, (0, 0, 0), 2)
        cv2.putText(img, "DOC format processing not fully implemented", (50, 200), 
                   cv2.FONT_HERSHEY_SIMPLEX, 0.8, (0, 0, 0), 1)
        cv2.putText(img, "Convert to DOCX for better processing", (50, 300), 
                   cv2.FONT_HERSHEY_SIMPLEX, 0.8, (0, 0, 0), 1)
        
        cv2.imwrite(str(output_path), img)
        return [str(output_path)]
    
    def _process_docx(self, file_path: str, output_dir: str) -> List[str]:
        """Process DOCX document.
        
        Args:
            file_path: Path to DOCX file
            output_dir: Output directory for images
            
        Returns:
            List of generated image paths
        """
        if not self.doc_available:
            raise RuntimeError("DOCX processing not available")
        
        try:
            import docx2txt
            
            # Extract text from DOCX
            text = docx2txt.process(file_path)
            
            # Create image representation
            output_path = Path(output_dir) / "page_000.png"
            
            # Create image with text content
            img = np.ones((1200, 800, 3), dtype=np.uint8) * 255
            
            # Add title
            cv2.putText(img, "DOCX Document", (50, 100), 
                       cv2.FONT_HERSHEY_SIMPLEX, 1.5, (0, 0, 0), 2)
            
            # Add text content (truncated for display)
            lines = text.split('\n')[:30]  # First 30 lines
            y_pos = 200
            for line in lines:
                if line.strip():
                    # Truncate long lines
                    display_line = line[:80] + "..." if len(line) > 80 else line
                    cv2.putText(img, display_line, (50, y_pos), 
                               cv2.FONT_HERSHEY_SIMPLEX, 0.6, (0, 0, 0), 1)
                    y_pos += 30
                    if y_pos > 1100:  # Prevent overflow
                        break
            
            cv2.imwrite(str(output_path), img)
            return [str(output_path)]
            
        except Exception as e:
            logger.error(f"Error processing DOCX {file_path}: {e}")
            raise
    
    def _process_ppt(self, file_path: str, output_dir: str) -> List[str]:
        """Process PPT document.
        
        Args:
            file_path: Path to PPT file
            output_dir: Output directory for images
            
        Returns:
            List of generated image paths
        """
        if not self.ppt_available:
            raise RuntimeError("PPT processing not available")
        
        # For now, create a placeholder image
        # In production, you would use python-pptx to extract content
        logger.warning("PPT processing is limited. Consider converting to PPTX first.")
        
        output_path = Path(output_dir) / "page_000.png"
        
        # Create a simple image with presentation info
        img = np.ones((800, 600, 3), dtype=np.uint8) * 255
        cv2.putText(img, f"PPT Presentation: {Path(file_path).name}", (50, 100), 
                   cv2.FONT_HERSHEY_SIMPLEX, 1, (0, 0, 0), 2)
        cv2.putText(img, "PPT format processing not fully implemented", (50, 200), 
                   cv2.FONT_HERSHEY_SIMPLEX, 0.8, (0, 0, 0), 1)
        cv2.putText(img, "Convert to PPTX for better processing", (50, 300), 
                   cv2.FONT_HERSHEY_SIMPLEX, 0.8, (0, 0, 0), 1)
        
        cv2.imwrite(str(output_path), img)
        return [str(output_path)]
    
    def _process_pptx(self, file_path: str, output_dir: str) -> List[str]:
        """Process PPTX document.
        
        Args:
            file_path: Path to PPTX file
            output_dir: Output directory for images
            
        Returns:
            List of generated image paths
        """
        if not self.ppt_available:
            raise RuntimeError("PPTX processing not available")
        
        try:
            from pptx import Presentation
            
            # Load presentation
            prs = Presentation(file_path)
            image_paths = []
            
            for i, slide in enumerate(prs.slides):
                # Create image representation of slide
                output_path = Path(output_dir) / f"slide_{i:03d}.png"
                
                # Create image with slide content
                img = np.ones((800, 600, 3), dtype=np.uint8) * 255
                
                # Add slide title
                cv2.putText(img, f"Slide {i+1}", (50, 100), 
                           cv2.FONT_HERSHEY_SIMPLEX, 1.2, (0, 0, 0), 2)
                
                # Extract text from slide
                y_pos = 200
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        # Truncate long text
                        text = shape.text[:60] + "..." if len(shape.text) > 60 else shape.text
                        cv2.putText(img, text, (50, y_pos), 
                                   cv2.FONT_HERSHEY_SIMPLEX, 0.6, (0, 0, 0), 1)
                        y_pos += 30
                        if y_pos > 700:  # Prevent overflow
                            break
                
                cv2.imwrite(str(output_path), img)
                image_paths.append(str(output_path))
            
            return image_paths
            
        except Exception as e:
            logger.error(f"Error processing PPTX {file_path}: {e}")
            raise
    
    def _process_image(self, file_path: str, output_dir: str) -> List[str]:
        """Process image file (copy to output directory).
        
        Args:
            file_path: Path to image file
            output_dir: Output directory
            
        Returns:
            List containing the copied image path
        """
        # For images, just copy to output directory
        input_path = Path(file_path)
        output_path = Path(output_dir) / input_path.name
        
        # Read and save image (this also validates the image)
        img = cv2.imread(file_path)
        if img is None:
            raise ValueError(f"Could not read image: {file_path}")
        
        cv2.imwrite(str(output_path), img)
        return [str(output_path)]
    
    def batch_process(self, file_paths: List[str], output_dir: str) -> Dict[str, List[str]]:
        """Process multiple documents in batch.
        
        Args:
            file_paths: List of file paths to process
            output_dir: Base output directory
            
        Returns:
            Dictionary mapping input files to output image lists
        """
        results = {}
        
        for file_path in file_paths:
            try:
                # Create subdirectory for each document
                doc_name = Path(file_path).stem
                doc_output_dir = Path(output_dir) / doc_name
                
                image_paths = self.process_document(file_path, str(doc_output_dir))
                results[file_path] = image_paths
                
            except Exception as e:
                logger.error(f"Failed to process {file_path}: {e}")
                results[file_path] = []
        
        return results

def process_document(file_path: str, output_dir: str = None) -> List[str]:
    """Convenience function to process a single document.
    
    Args:
        file_path: Path to input document
        output_dir: Output directory for images
        
    Returns:
        List of generated image paths
    """
    processor = DocumentProcessor()
    return processor.process_document(file_path, output_dir)

def batch_process_documents(file_paths: List[str], output_dir: str) -> Dict[str, List[str]]:
    """Convenience function to process multiple documents.
    
    Args:
        file_paths: List of file paths to process
        output_dir: Base output directory
        
    Returns:
        Dictionary mapping input files to output image lists
    """
    processor = DocumentProcessor()
    return processor.batch_process(file_paths, output_dir)

if __name__ == "__main__":
    import argparse
    
    parser = argparse.ArgumentParser(description="Process documents to images")
    parser.add_argument('--input', required=True, help='Input file or directory')
    parser.add_argument('--output', required=True, help='Output directory')
    
    args = parser.parse_args()
    
    input_path = Path(args.input)
    
    if input_path.is_file():
        # Process single file
        image_paths = process_document(str(input_path), args.output)
        print(f"Processed {input_path} -> {len(image_paths)} images")
        for img_path in image_paths:
            print(f"  {img_path}")
    
    elif input_path.is_dir():
        # Process directory
        file_paths = []
        for fmt in ['.pdf', '.doc', '.docx', '.ppt', '.pptx', '.png', '.jpg', '.jpeg']:
            file_paths.extend(input_path.glob(f"*{fmt}"))
            file_paths.extend(input_path.glob(f"*{fmt.upper()}"))
        
        if file_paths:
            results = batch_process_documents([str(p) for p in file_paths], args.output)
            print(f"Processed {len(file_paths)} files:")
            for file_path, image_paths in results.items():
                print(f"  {file_path} -> {len(image_paths)} images")
        else:
            print("No supported files found in directory")
    
    else:
        print(f"Input path does not exist: {input_path}")
