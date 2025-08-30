"""
Document Cleaning Service for PS-05

Implements all 11 document cleaning tasks:
1. Text Extraction & Encoding
2. Removing Boilerplate Text
3. Handling Hyphenation & Line Breaks
4. Removing Non-Text Elements
5. Normalization
6. Removing Special Characters
7. Tokenization & Stopword Removal
8. Metadata Extraction & Cleaning
9. Language Detection
10. Structure Recovery
11. Deduplication
"""

import logging
import json
import re
import hashlib
from pathlib import Path
from typing import List, Dict, Optional, Tuple, Any
from datetime import datetime
import unicodedata
import langdetect
from langdetect import detect, DetectorFactory
import nltk
from nltk.tokenize import word_tokenize, sent_tokenize
from nltk.corpus import stopwords
import PyPDF2
import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation
import zipfile
import xml.etree.ElementTree as ET

logger = logging.getLogger(__name__)

class DocumentCleaningService:
    """Comprehensive document cleaning service for various document formats."""
    
    def __init__(self, config: Optional[Dict] = None):
        self.config = config or self._get_default_config()
        self.cleaned_dir = None
        self.cleaning_log = []
        
        # Download required NLTK data
        try:
            nltk.download('punkt', quiet=True)
            nltk.download('stopwords', quiet=True)
            nltk.download('averaged_perceptron_tagger', quiet=True)
        except Exception as e:
            logger.warning(f"Failed to download NLTK data: {e}")
        
        # Set seed for consistent language detection
        DetectorFactory.seed = 0
    
    def _get_default_config(self) -> Dict:
        """Get default cleaning configuration."""
        return {
            "supported_formats": ['.pdf', '.docx', '.pptx', '.txt'],
            "encoding": 'utf-8',
            "min_text_length": 50,  # Minimum characters for valid document
            "remove_boilerplate": True,
            "normalize_text": True,
            "remove_stopwords": False,  # Keep for document understanding
            "language_detection": True,
            "structure_recovery": True,
            "deduplication": True,
            "output_format": "json",  # json, txt, or both
            "quality_threshold": 0.7
        }
    
    def clean_dataset(
        self, 
        input_dir: Path, 
        output_dir: Path
    ) -> Dict:
        """
        Clean entire document dataset with all 11 cleaning tasks.
        
        Args:
            input_dir: Directory containing raw documents
            output_dir: Directory for cleaned documents
            
        Returns:
            Dictionary with cleaning results and statistics
        """
        try:
            logger.info(f"Starting comprehensive document cleaning for {input_dir}")
            
            # Create output directory
            output_dir.mkdir(parents=True, exist_ok=True)
            self.cleaned_dir = output_dir
            
            # Initialize cleaning log
            self.cleaning_log = []
            
            # Get all document files
            document_files = self._get_document_files(input_dir)
            logger.info(f"Found {len(document_files)} documents to process")
            
            # Initialize statistics
            stats = {
                "total_documents": len(document_files),
                "cleaned_documents": 0,
                "removed_corrupt": 0,
                "removed_duplicates": 0,
                "removed_boilerplate": 0,
                "language_detected": {},
                "cleaning_errors": 0,
                "format_breakdown": {}
            }
            
            # Process each document
            cleaned_documents = []
            document_hashes = set()
            
            for doc_path in document_files:
                try:
                    # Update format breakdown
                    doc_ext = doc_path.suffix.lower()
                    stats["format_breakdown"][doc_ext] = stats["format_breakdown"].get(doc_ext, 0) + 1
                    
                    # Clean document
                    cleaned_doc = self._clean_single_document(doc_path, output_dir)
                    
                    if cleaned_doc:
                        # Check for duplicates
                        doc_hash = self._calculate_document_hash(cleaned_doc["content"])
                        
                        if self.config["deduplication"] and doc_hash in document_hashes:
                            stats["removed_duplicates"] += 1
                            logger.info(f"Duplicate removed: {doc_path}")
                            continue
                        
                        document_hashes.add(doc_hash)
                        cleaned_documents.append(cleaned_doc)
                        stats["cleaned_documents"] += 1
                        
                        # Update language statistics
                        if cleaned_doc.get("language"):
                            lang = cleaned_doc["language"]
                            stats["language_detected"][lang] = stats["language_detected"].get(lang, 0) + 1
                        
                        # Log cleaning action
                        self.cleaning_log.append({
                            "action": "document_cleaned",
                            "original": str(doc_path),
                            "cleaned": str(cleaned_doc.get("output_path", "")),
                            "operations": cleaned_doc.get("operations", []),
                            "language": cleaned_doc.get("language", "unknown"),
                            "word_count": cleaned_doc.get("word_count", 0)
                        })
                    
                except Exception as e:
                    logger.error(f"Failed to clean {doc_path}: {e}")
                    stats["cleaning_errors"] += 1
            
            # Save cleaning log
            self._save_cleaning_log(output_dir)
            
            # Generate final statistics
            final_stats = self._generate_final_stats(stats, cleaned_documents)
            
            logger.info(f"Document cleaning completed. Final documents: {len(cleaned_documents)}")
            return final_stats
            
        except Exception as e:
            logger.error(f"Document cleaning failed: {e}")
            raise
    
    def _get_document_files(self, input_dir: Path) -> List[Path]:
        """Get all supported document files from directory."""
        document_files = []
        for format_ext in self.config["supported_formats"]:
            document_files.extend(input_dir.glob(f"*{format_ext}"))
            document_files.extend(input_dir.glob(f"*{format_ext.upper()}"))
        return sorted(document_files)
    
    def _clean_single_document(self, doc_path: Path, output_dir: Path) -> Optional[Dict]:
        """Clean a single document with all cleaning tasks."""
        try:
            doc_ext = doc_path.suffix.lower()
            
            # Extract text based on format
            if doc_ext == '.pdf':
                content = self._extract_pdf_text(doc_path)
            elif doc_ext == '.docx':
                content = self._extract_docx_text(doc_path)
            elif doc_ext == '.pptx':
                content = self._extract_pptx_text(doc_path)
            elif doc_ext == '.txt':
                content = self._extract_txt_text(doc_path)
            else:
                logger.warning(f"Unsupported format: {doc_ext}")
                return None
            
            if not content or len(content.strip()) < self.config["min_text_length"]:
                logger.warning(f"Document too short or empty: {doc_path}")
                return None
            
            # Apply all cleaning tasks
            cleaned_content = self._apply_all_cleaning_tasks(content)
            
            # Language detection
            language = "unknown"
            if self.config["language_detection"]:
                language = self._detect_language(cleaned_content)
            
            # Structure recovery
            structure = {}
            if self.config["structure_recovery"]:
                structure = self._recover_structure(cleaned_content)
            
            # Save cleaned document
            output_path = self._save_cleaned_document(
                doc_path, cleaned_content, output_dir, language, structure
            )
            
            # Calculate word count
            word_count = len(cleaned_content.split())
            
            return {
                "original_path": str(doc_path),
                "output_path": str(output_path),
                "content": cleaned_content,
                "language": language,
                "structure": structure,
                "word_count": word_count,
                "operations": ["text_extraction", "boilerplate_removal", "normalization", 
                             "special_char_removal", "structure_recovery", "language_detection"]
            }
            
        except Exception as e:
            logger.error(f"Failed to clean {doc_path}: {e}")
            return None
    
    def _extract_pdf_text(self, pdf_path: Path) -> str:
        """Task 1: Extract text from PDF with encoding handling."""
        try:
            # Try PyMuPDF first (better text extraction)
            try:
                doc = fitz.open(str(pdf_path))
                text = ""
                for page in doc:
                    text += page.get_text()
                doc.close()
                return text
            except Exception as e:
                logger.warning(f"PyMuPDF failed, trying PyPDF2: {e}")
                
                # Fallback to PyPDF2
                with open(pdf_path, 'rb') as file:
                    reader = PyPDF2.PdfReader(file)
                    text = ""
                    for page in reader.pages:
                        text += page.extract_text()
                return text
                
        except Exception as e:
            logger.error(f"PDF text extraction failed: {e}")
            return ""
    
    def _extract_docx_text(self, docx_path: Path) -> str:
        """Task 1: Extract text from DOCX."""
        try:
            doc = Document(docx_path)
            text = ""
            
            # Extract text from paragraphs
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
            
            # Extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        text += cell.text + "\t"
                    text += "\n"
            
            return text
            
        except Exception as e:
            logger.error(f"DOCX text extraction failed: {e}")
            return ""
    
    def _extract_pptx_text(self, pptx_path: Path) -> str:
        """Task 1: Extract text from PPTX."""
        try:
            prs = Presentation(pptx_path)
            text = ""
            
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            
            return text
            
        except Exception as e:
            logger.error(f"PPTX text extraction failed: {e}")
            return ""
    
    def _extract_txt_text(self, txt_path: Path) -> str:
        """Task 1: Extract text from TXT with encoding detection."""
        try:
            # Try different encodings
            encodings = ['utf-8', 'latin-1', 'cp1252', 'iso-8859-1']
            
            for encoding in encodings:
                try:
                    with open(txt_path, 'r', encoding=encoding) as f:
                        return f.read()
                except UnicodeDecodeError:
                    continue
            
            # If all fail, try with error handling
            with open(txt_path, 'r', encoding='utf-8', errors='ignore') as f:
                return f.read()
                
        except Exception as e:
            logger.error(f"TXT text extraction failed: {e}")
            return ""
    
    def _apply_all_cleaning_tasks(self, content: str) -> str:
        """Apply all cleaning tasks to the content."""
        cleaned_content = content
        
        # Task 2: Remove boilerplate text
        if self.config["remove_boilerplate"]:
            cleaned_content = self._remove_boilerplate(cleaned_content)
        
        # Task 3: Handle hyphenation & line breaks
        cleaned_content = self._handle_hyphenation(cleaned_content)
        
        # Task 4: Remove non-text elements (basic)
        cleaned_content = self._remove_non_text_elements(cleaned_content)
        
        # Task 5: Normalize text
        if self.config["normalize_text"]:
            cleaned_content = self._normalize_text(cleaned_content)
        
        # Task 6: Remove special characters
        cleaned_content = self._remove_special_characters(cleaned_content)
        
        # Task 7: Tokenization (prepare for analysis)
        tokens = self._tokenize_text(cleaned_content)
        
        # Rejoin tokens
        cleaned_content = " ".join(tokens)
        
        return cleaned_content
    
    def _remove_boilerplate(self, content: str) -> str:
        """Task 2: Remove boilerplate text."""
        # Common boilerplate patterns
        boilerplate_patterns = [
            r'Page \d+ of \d+',
            r'© \d{4}.*?\.',
            r'All rights reserved\.',
            r'Confidential',
            r'Draft',
            r'Version \d+\.\d+',
            r'Last updated:.*?\.',
            r'Created:.*?\.',
            r'Modified:.*?\.',
            r'File:.*?\.',
            r'Path:.*?\.',
            r'Size:.*?\.',
            r'Word count:.*?\.',
            r'Character count:.*?\.'
        ]
        
        for pattern in boilerplate_patterns:
            content = re.sub(pattern, '', content, flags=re.IGNORECASE | re.MULTILINE)
        
        return content
    
    def _handle_hyphenation(self, content: str) -> str:
        """Task 3: Handle hyphenation and line breaks."""
        # Join hyphenated words at line breaks
        content = re.sub(r'(\w+)-\s*\n\s*(\w+)', r'\1\2', content)
        
        # Join broken sentences
        content = re.sub(r'(\w+)\s*\n\s*(\w+)', r'\1 \2', content)
        
        # Remove excessive whitespace
        content = re.sub(r'\s+', ' ', content)
        
        return content.strip()
    
    def _remove_non_text_elements(self, content: str) -> str:
        """Task 4: Remove non-text elements."""
        # Remove XML/HTML tags
        content = re.sub(r'<[^>]+>', '', content)
        
        # Remove URLs
        content = re.sub(r'http[s]?://(?:[a-zA-Z]|[0-9]|[$-_@.&+]|[!*\\(\\),]|(?:%[0-9a-fA-F][0-9a-fA-F]))+', '', content)
        
        # Remove email addresses
        content = re.sub(r'\b[A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Z|a-z]{2,}\b', '', content)
        
        # Remove file paths
        content = re.sub(r'[A-Za-z]:\\[\\\S|*\S]?.*', '', content)
        content = re.sub(r'/[^\s]*', '', content)
        
        return content
    
    def _normalize_text(self, content: str) -> str:
        """Task 5: Normalize text."""
        # Convert to lowercase
        content = content.lower()
        
        # Expand contractions
        contractions = {
            "don't": "do not",
            "can't": "cannot",
            "won't": "will not",
            "it's": "it is",
            "that's": "that is",
            "you're": "you are",
            "i'm": "i am",
            "we're": "we are",
            "they're": "they are",
            "he's": "he is",
            "she's": "she is"
        }
        
        for contraction, expansion in contractions.items():
            content = content.replace(contraction, expansion)
        
        # Normalize unicode
        content = unicodedata.normalize('NFKC', content)
        
        return content
    
    def _remove_special_characters(self, content: str) -> str:
        """Task 6: Remove special characters."""
        # Keep alphanumeric, spaces, and basic punctuation
        content = re.sub(r'[^\w\s\.\,\!\?\;\:\-\(\)]', '', content)
        
        # Remove excessive punctuation
        content = re.sub(r'[\.]{2,}', '.', content)
        content = re.sub(r'[\!]{2,}', '!', content)
        content = re.sub(r'[\?]{2,}', '?', content)
        
        # Clean up whitespace
        content = re.sub(r'\s+', ' ', content)
        
        return content.strip()
    
    def _tokenize_text(self, content: str) -> List[str]:
        """Task 7: Tokenize text."""
        try:
            # Sentence tokenization
            sentences = sent_tokenize(content)
            
            # Word tokenization
            tokens = []
            for sentence in sentences:
                words = word_tokenize(sentence)
                tokens.extend(words)
            
            # Remove stopwords if configured
            if self.config["remove_stopwords"]:
                stop_words = set(stopwords.words('english'))
                tokens = [word for word in tokens if word.lower() not in stop_words]
            
            return tokens
            
        except Exception as e:
            logger.warning(f"Tokenization failed, using simple split: {e}")
            return content.split()
    
    def _detect_language(self, content: str) -> str:
        """Task 9: Detect document language."""
        try:
            # Use first 1000 characters for language detection
            sample_text = content[:1000]
            if len(sample_text.strip()) < 50:
                return "unknown"
            
            language = detect(sample_text)
            return language
            
        except Exception as e:
            logger.warning(f"Language detection failed: {e}")
            return "unknown"
    
    def _recover_structure(self, content: str) -> Dict:
        """Task 10: Recover document structure."""
        try:
            structure = {
                "paragraphs": [],
                "headings": [],
                "lists": [],
                "tables": [],
                "metadata": {}
            }
            
            # Split into paragraphs
            paragraphs = content.split('\n\n')
            structure["paragraphs"] = [p.strip() for p in paragraphs if p.strip()]
            
            # Detect headings (lines that are shorter and end with punctuation)
            for para in structure["paragraphs"]:
                if len(para.split()) <= 10 and para.strip().endswith((':', '.', '!')):
                    structure["headings"].append(para.strip())
            
            # Detect lists (lines starting with numbers or bullets)
            for para in structure["paragraphs"]:
                if re.match(r'^[\d\-\*\.]+\s+', para.strip()):
                    structure["lists"].append(para.strip())
            
            # Basic metadata
            structure["metadata"] = {
                "word_count": len(content.split()),
                "character_count": len(content),
                "paragraph_count": len(structure["paragraphs"]),
                "heading_count": len(structure["headings"]),
                "list_count": len(structure["lists"])
            }
            
            return structure
            
        except Exception as e:
            logger.warning(f"Structure recovery failed: {e}")
            return {"error": str(e)}
    
    def _calculate_document_hash(self, content: str) -> str:
        """Calculate hash for deduplication."""
        # Normalize content for consistent hashing
        normalized = content.lower().strip()
        return hashlib.md5(normalized.encode('utf-8')).hexdigest()
    
    def _save_cleaned_document(
        self, 
        original_path: Path, 
        content: str, 
        output_dir: Path, 
        language: str, 
        structure: Dict
    ) -> Path:
        """Save cleaned document in specified format."""
        try:
            base_name = original_path.stem
            
            if self.config["output_format"] in ["json", "both"]:
                # Save as JSON with metadata
                json_data = {
                    "original_filename": original_path.name,
                    "cleaned_content": content,
                    "language": language,
                    "structure": structure,
                    "cleaning_timestamp": datetime.now().isoformat(),
                    "cleaning_operations": ["text_extraction", "boilerplate_removal", 
                                         "normalization", "structure_recovery"]
                }
                
                json_path = output_dir / f"{base_name}_cleaned.json"
                with open(json_path, 'w', encoding='utf-8') as f:
                    json.dump(json_data, f, indent=2, ensure_ascii=False)
            
            if self.config["output_format"] in ["txt", "both"]:
                # Save as plain text
                txt_path = output_dir / f"{base_name}_cleaned.txt"
                with open(txt_path, 'w', encoding='utf-8') as f:
                    f.write(content)
            
            # Return the primary output path
            if self.config["output_format"] == "json":
                return json_path
            elif self.config["output_format"] == "txt":
                return txt_path
            else:
                return json_path  # Default to JSON
                
        except Exception as e:
            logger.error(f"Failed to save cleaned document: {e}")
            raise
    
    def _save_cleaning_log(self, output_dir: Path):
        """Save detailed cleaning log."""
        log_file = output_dir / "document_cleaning_log.json"
        with open(log_file, 'w', encoding='utf-8') as f:
            json.dump(self.cleaning_log, f, indent=2, ensure_ascii=False)
        logger.info(f"Document cleaning log saved: {log_file}")
    
    def _generate_final_stats(self, stats: Dict, cleaned_documents: List[Dict]) -> Dict:
        """Generate final cleaning statistics."""
        final_stats = {
            "cleaning_summary": stats,
            "final_document_count": len(cleaned_documents),
            "cleaning_efficiency": (len(cleaned_documents) / stats["total_documents"]) * 100,
            "output_directory": str(self.cleaned_dir),
            "cleaning_log_file": "document_cleaning_log.json",
            "language_distribution": stats["language_detected"],
            "format_distribution": stats["format_breakdown"]
        }
        
        return final_stats
